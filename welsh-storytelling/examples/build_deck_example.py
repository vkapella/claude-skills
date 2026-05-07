"""
Demo: BNY RAP Sites deck reworked through Michael Welsh's storytelling frameworks.

Applies:
  - 5 Slide Rule (reduced front-of-deck to 6 business slides)
  - Story Arc: Setup -> What -> How -> Why -> Magic Wand
  - Rule of 3 on every content slide
  - Diamond facets on the core message
  - Ends with wanting (clear ask)

Original deck had 31 slides with a flat opening that buried the "why."
This rework focuses on the BUSINESS opening that the COO needs.
The technical appendix (slides 3-5, 9-31 from original) stays intact
and would follow these 6 reworked slides.
"""

import sys, os
sys.path.insert(0, '/mnt/skills/organization/wwt-presentation')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================================
# WWT Brand Constants
# ============================================================================
WWT_DARK_BLUE  = RGBColor(0x1C, 0x00, 0x87)
WWT_LIGHT_BLUE = RGBColor(0x00, 0x86, 0xEA)
WWT_RED        = RGBColor(0xEE, 0x28, 0x2A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
BODY_GRAY      = RGBColor(0x26, 0x26, 0x26)
DARK_BLUE_75   = RGBColor(0x49, 0x33, 0x9F)
LIGHT_BLUE_75  = RGBColor(0x33, 0x9E, 0xEE)
LIGHT_BLUE_50  = RGBColor(0x66, 0xB6, 0xF2)
LIGHT_BLUE_TINT = RGBColor(0xE6, 0xF3, 0xFD)
NAVY           = RGBColor(0x1D, 0x1E, 0x48)
ORANGE         = RGBColor(0xFB, 0x55, 0x0E)
LIGHT_GRAY_BG  = RGBColor(0xF5, 0xF7, 0xFA)
MID_GRAY       = RGBColor(0x6B, 0x72, 0x80)

ASSETS_DIR = "/mnt/skills/organization/wwt-presentation/assets"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CONTENT_LEFT  = Inches(0.8)
CONTENT_WIDTH = Inches(11.733)
TITLE_TOP     = Inches(0.6)

# ============================================================================
# Helpers
# ============================================================================
def set_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bg_image(slide, image_path):
    pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    sp = pic._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)
    return pic

def add_gradient_line(slide):
    height = Inches(0.05)
    y = SLIDE_H - height
    return slide.shapes.add_picture(os.path.join(ASSETS_DIR, "gradient_line.png"),
                                     Inches(0), y, SLIDE_W, height)

def add_monogram(slide):
    mono_w = Inches(0.45)
    mono_h = Inches(0.45)
    x = SLIDE_W - mono_w - Inches(0.4)
    y = SLIDE_H - mono_h - Inches(0.25)
    return slide.shapes.add_picture(os.path.join(ASSETS_DIR, "monogram_color.png"),
                                     x, y, mono_w, mono_h)

def add_accent_bar_top(slide, color=None):
    if color is None: color = WWT_LIGHT_BLUE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_textbox(slide, left, top, width, height, text, font_size=Pt(20),
                color=BODY_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial", anchor=MSO_ANCHOR.TOP, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    bodyPr = tf.paragraphs[0]._p.getparent().find(qn('a:bodyPr'))
    if bodyPr is not None:
        if anchor == MSO_ANCHOR.MIDDLE:
            bodyPr.set('anchor', 'ctr')
        elif anchor == MSO_ANCHOR.BOTTOM:
            bodyPr.set('anchor', 'b')

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_block(slide, title_text, eyebrow_text=""):
    """Standard slide title with optional eyebrow label above."""
    if eyebrow_text:
        add_textbox(slide, CONTENT_LEFT, Inches(0.45), CONTENT_WIDTH, Inches(0.35),
                    eyebrow_text.upper(), font_size=Pt(11), color=ORANGE, bold=True)
        add_textbox(slide, CONTENT_LEFT, Inches(0.75), CONTENT_WIDTH, Inches(0.85),
                    title_text, font_size=Pt(34), color=WWT_DARK_BLUE, bold=True)
    else:
        add_textbox(slide, CONTENT_LEFT, TITLE_TOP, CONTENT_WIDTH, Inches(1.0),
                    title_text, font_size=Pt(36), color=WWT_LIGHT_BLUE, bold=True)
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, CONTENT_LEFT, Inches(1.55), Inches(2.0), Inches(0.04)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = WWT_LIGHT_BLUE
    sep.line.fill.background()

def add_speaker_notes(slide, notes_text):
    """Attach speaker notes (the BME of the slide for the presenter)."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

# ============================================================================
# Build the deck
# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# SLIDE 1 — TITLE (sets the magic-wand promise)
# Welsh: A title slide should make a promise the deck keeps.
# Original promised "Intelligent Lifecycle Approach" but delivered a deployment plan.
# Reworked title is honest about what's being decided.
# ----------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_bg_image(slide, os.path.join(ASSETS_DIR, "bg_navy_graphic_device.png"))

add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.5),
            "BNY  ·  RAP MICROBRANCH PROGRAM",
            font_size=Pt(14), color=LIGHT_BLUE_75, bold=True)

add_textbox(slide, Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.4),
            "11 sites. 6 months. One decision.",
            font_size=Pt(48), color=WHITE, bold=True)

accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.6), Inches(2.0), Inches(0.05)
)
accent.fill.solid()
accent.fill.fore_color.rgb = WWT_RED
accent.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.7),
            "Choosing the path to a stable, segmented, scalable branch network",
            font_size=Pt(22), color=LIGHT_BLUE_75)

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(6.0), Inches(0.4),
            "WWT  ·  Prepared for BNY COO  ·  March 2026",
            font_size=Pt(13), color=RGBColor(0xCC, 0xCC, 0xCC))

add_speaker_notes(slide, """SETUP — Open the deck.

BME of this slide: Begin = a real beginning, not a logo. Middle = the audience grasps the scope (11 sites, 6 months) before any technical detail. End = they know they are here to make ONE decision today.

What to say (verbatim opener, 30 seconds):
"Thanks for the time. This deck has one purpose — to help us land on the right network for the eleven new RAP sites we're standing up over the next six months. There's a recommendation at the end. Before that, I want to make sure we agree on what broke last time, because that's what's shaping the choice."

PAUSE three seconds before clicking. Let the room settle.""")

# ----------------------------------------------------------------------------
# SLIDE 2 — SETUP (Why we're here)
# Welsh: The original deck never had a real Setup. The "Why Microbranch" slide
# was buried at slide 11. We promote it to slide 2 — this is the BMEs of the
# situation that triggered this whole program.
# Rule of 3: Three failures, not four (we merge "fragmented management" into
# the broader story of operational pain).
# ----------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, WHITE)
add_accent_bar_top(slide)

add_title_block(slide, "The last build broke. Here's how.", eyebrow_text="Setup  ·  Why we're here")

# Three failure cards (Rule of 3)
failures = [
    ("Connectivity\nwas brittle",
     "Aruba RAPs dropped IPs on every VPN reconnect, breaking static-IP services and degrading user experience. Internet-only meant any ISP outage took the branch dark."),
    ("Traffic\nwas blind",
     "Voice, video, and trading traffic competed for the same WAN path. No QoS, no application-aware routing, no way to protect what mattered most."),
    ("Operations\nwere reactive",
     "No single dashboard. No streaming telemetry. The NOC found out about faults the same way users did — by watching them happen.")
]

card_w = Inches(3.85)
card_h = Inches(4.35)
card_y = Inches(2.1)
gap = Inches(0.18)
start_x = Inches(0.8)

for i, (heading, body) in enumerate(failures):
    x = start_x + (card_w + gap) * i

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY_BG
    card.line.color.rgb = LIGHT_BLUE_TINT
    card.line.width = Pt(0.5)

    # Red accent bar at top of card
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.08))
    acc.fill.solid()
    acc.fill.fore_color.rgb = WWT_RED
    acc.line.fill.background()

    # Big number
    add_textbox(slide, x + Inches(0.3), card_y + Inches(0.3), Inches(1.0), Inches(0.8),
                f"0{i+1}", font_size=Pt(40), color=WWT_LIGHT_BLUE, bold=True)

    # Heading
    add_textbox(slide, x + Inches(0.3), card_y + Inches(1.15), card_w - Inches(0.6), Inches(1.2),
                heading, font_size=Pt(22), color=WWT_DARK_BLUE, bold=True)

    # Body
    add_textbox(slide, x + Inches(0.3), card_y + Inches(2.4), card_w - Inches(0.6), Inches(1.85),
                body, font_size=Pt(13), color=BODY_GRAY)

# Bottom takeaway
add_textbox(slide, CONTENT_LEFT, Inches(6.65), CONTENT_WIDTH, Inches(0.4),
            "Three root causes — well-documented by BNY. The new design has to fix all three, or we'll be having this meeting again in eighteen months.",
            font_size=Pt(13), color=MID_GRAY, italic=True)

add_monogram(slide)
add_gradient_line(slide)

add_speaker_notes(slide, """SETUP (continued) — Establish the BMEs of the problem.

BME: Begin = the previous build broke. Middle = three specific, documented failures. End = these three failures are the design constraints — not WWT's opinions.

What to say (90 seconds):
Walk through each card briefly. Use the Rule of 3 — speak to all three, but linger on the one that matters most to your audience. For a COO, that's almost always #3 (operations were reactive) — frame it as: "your team was paid to be proactive, but the platform forced them to be reactive."

REPETITION WITH VARIATION: Come back to "fix all three or we'll be back here in eighteen months" later in the deck. This is your message-discipline anchor.

WHAT NOT TO SAY: Do not blame Aruba. The failures are documented; the COO knows. Naming a vendor as the villain makes you look small.""")

# ----------------------------------------------------------------------------
# SLIDE 3 — WHAT (The core message — the Diamond's top facet)
# Welsh: One slide. One message. This is the placemat slide.
# Original deck had 5 separate boxes on the overview. We collapse to ONE statement.
# ----------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, WHITE)
add_accent_bar_top(slide)

add_title_block(slide, "What we're proposing", eyebrow_text="What  ·  The core message")

# Big anchor statement (the Diamond's Point)
quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    CONTENT_LEFT, Inches(2.0), CONTENT_WIDTH, Inches(2.3))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = NAVY
quote_box.line.fill.background()

add_textbox(slide, CONTENT_LEFT + Inches(0.5), Inches(2.25), CONTENT_WIDTH - Inches(1.0), Inches(1.8),
            "Deploy 11 microbranches and 4 head-ends in 6 months on Cisco SD-WAN — the only option that fixes all three root causes and aligns with the broader Branch Refresh.",
            font_size=Pt(26), color=WHITE, bold=False, anchor=MSO_ANCHOR.MIDDLE)

# Small Red dot
dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
    CONTENT_LEFT + Inches(0.5), Inches(2.4), Inches(0.15), Inches(0.15))
dot.fill.solid()
dot.fill.fore_color.rgb = WWT_RED
dot.line.fill.background()

# Three supporting facets (Rule of 3 — the Diamond's other facets)
add_textbox(slide, CONTENT_LEFT, Inches(4.7), CONTENT_WIDTH, Inches(0.4),
            "WHY THIS, WHY NOW",
            font_size=Pt(11), color=ORANGE, bold=True)

facets = [
    ("Forward-aligned", "Same architecture as the broader Branch Refresh. No throwaway investment."),
    ("Fully compliant", "Meets all 12 design requirements — including end-to-end segmentation."),
    ("Time-boxed", "First 11 sites turnkey by end of Q3. Pilot for the WWT Branch-in-a-Box program.")
]

facet_w = Inches(3.85)
facet_y = Inches(5.15)
for i, (h, b) in enumerate(facets):
    x = CONTENT_LEFT + (facet_w + Inches(0.18)) * i

    # Subtle left border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, facet_y, Inches(0.04), Inches(1.4))
    border.fill.solid()
    border.fill.fore_color.rgb = WWT_LIGHT_BLUE
    border.line.fill.background()

    add_textbox(slide, x + Inches(0.2), facet_y, facet_w - Inches(0.2), Inches(0.4),
                h, font_size=Pt(15), color=WWT_DARK_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.2), facet_y + Inches(0.4), facet_w - Inches(0.2), Inches(1.0),
                b, font_size=Pt(12), color=BODY_GRAY)

add_monogram(slide)
add_gradient_line(slide)

add_speaker_notes(slide, """WHAT — The core message. The Diamond's top facet.

This is the ONE thing the COO needs to remember if they remember nothing else.

What to say (60 seconds):
"Here's the recommendation in one sentence — read it aloud — and here's the why in three short reasons. We'll go deeper on each in the next slides."

THE DIAMOND IN ACTION:
- The Point (top facet): the navy box statement.
- Context (side facets): the three supporting reasons.
- Evidence (middle facets): the technical conformance + cost slides that follow.
- Application (bottom facets): what this means for BNY operationally.
- Call to Action: comes on the closing slide.

REPETITION WITH VARIATION — three ways to land this same message later:
1. "One architecture, branch to data center" (technical framing)
2. "Buy once, use twice" (financial framing)
3. "The pilot for everything that comes after" (strategic framing)
Pick the one that fits the question being asked.""")

# ----------------------------------------------------------------------------
# SLIDE 4 — HOW (The proof — Rule of 3 on the technical comparison)
# Welsh: 5 things on a slide? Pick three. Speak to those three.
# Original conformance table had 12 rows. We surface the 3 that decide it.
# ----------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, WHITE)
add_accent_bar_top(slide)

add_title_block(slide, "Why Cisco SD-WAN, not Meraki", eyebrow_text="How  ·  The technical decision")

# Subtitle / framing
add_textbox(slide, CONTENT_LEFT, Inches(1.7), CONTENT_WIDTH, Inches(0.4),
            "Both options work. Three requirements separate them — and all three matter for the larger Branch Refresh.",
            font_size=Pt(15), color=MID_GRAY, italic=True)

# Three-column comparison
header_y = Inches(2.4)
row_y = Inches(3.05)

# Column headers
col_x = [Inches(0.8), Inches(5.5), Inches(9.4)]
col_w = [Inches(4.5), Inches(3.7), Inches(3.0)]

# Header row
hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), header_y, Inches(11.7), Inches(0.5))
hdr.fill.solid()
hdr.fill.fore_color.rgb = NAVY
hdr.line.fill.background()

add_textbox(slide, col_x[0] + Inches(0.2), header_y, col_w[0], Inches(0.5),
            "Requirement", font_size=Pt(13), color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(slide, col_x[1], header_y, col_w[1], Inches(0.5),
            "Cisco SD-WAN", font_size=Pt(13), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(slide, col_x[2], header_y, col_w[2], Inches(0.5),
            "Meraki", font_size=Pt(13), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Three rows (Rule of 3 on the 12-requirement table)
rows = [
    ("End-to-end segmentation\n(branch ↔ data center)",
     "Full SGT support",  True,
     "No SGT support",    False),
    ("Identity-based policy from NAC\nenforced at the WAN edge",
     "Native ISE integration", True,
     "Cannot enforce", False),
    ("Forward alignment with the\nbroader Branch Refresh",
     "Same architecture", True,
     "Separate platform", False),
]

for i, (req, sd_text, sd_ok, mer_text, mer_ok) in enumerate(rows):
    y = row_y + Inches(0.95) * i
    row_h = Inches(0.85)

    # Row background (alternating)
    row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(11.7), row_h)
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = LIGHT_BLUE_TINT if i % 2 == 0 else WHITE
    row_bg.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    row_bg.line.width = Pt(0.5)

    # Requirement
    add_textbox(slide, col_x[0] + Inches(0.2), y, col_w[0], row_h,
                req, font_size=Pt(13), color=BODY_GRAY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    # Cisco SD-WAN cell
    sd_color = RGBColor(0x05, 0x88, 0x4A) if sd_ok else WWT_RED
    sd_mark = "✓" if sd_ok else "✕"
    add_textbox(slide, col_x[1], y, Inches(0.5), row_h,
                sd_mark, font_size=Pt(22), color=sd_color, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, col_x[1] + Inches(0.5), y, col_w[1] - Inches(0.5), row_h,
                sd_text, font_size=Pt(13), color=BODY_GRAY, anchor=MSO_ANCHOR.MIDDLE)

    # Meraki cell
    mer_color = RGBColor(0x05, 0x88, 0x4A) if mer_ok else WWT_RED
    mer_mark = "✓" if mer_ok else "✕"
    add_textbox(slide, col_x[2], y, Inches(0.5), row_h,
                mer_mark, font_size=Pt(22), color=mer_color, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, col_x[2] + Inches(0.5), y, col_w[2] - Inches(0.5), row_h,
                mer_text, font_size=Pt(13), color=BODY_GRAY, anchor=MSO_ANCHOR.MIDDLE)

# Bottom callout
add_textbox(slide, CONTENT_LEFT, Inches(6.4), CONTENT_WIDTH, Inches(0.5),
            "Meraki is 50% cheaper and ships faster — but loses the segmentation, identity, and architectural alignment that the broader Refresh depends on. Full 12-requirement table is in the appendix.",
            font_size=Pt(12), color=MID_GRAY, italic=True)

add_monogram(slide)
add_gradient_line(slide)

add_speaker_notes(slide, """HOW — The technical proof.

BME: Begin = both options work (concession). Middle = three specific gaps that matter. End = price/speed don't override architecture for an enterprise refresh.

What to say (90 seconds):
"There are twelve design requirements. Cisco SD-WAN meets all twelve fully. Meraki meets nine fully and three partially or not at all. The full table is in the appendix. I want to spend our time on the three that decided this."

Walk each row in 20 seconds. Don't repeat the words on the slide — extend them.

POWER OF "I DON'T KNOW": If they ask about a specific requirement not on this slide, say: "Let me pull up the full table — I want to give you the right answer, not a fast one." Then go to the appendix slide.

REPETITION WITH VARIATION — three ways to frame the gap:
1. "Meraki is a great microbranch box. But this is bigger than microbranch."
2. "Cheaper now, more expensive later — because we'd build it twice."
3. "Meraki gets us through the door. SD-WAN gets us through the next decade."

WATCH FOR: If the COO leans in on cost (the $1.4M delta), do not defend. Acknowledge, then pivot to slide 5.""")

# ----------------------------------------------------------------------------
# SLIDE 5 — WHY (Cost framed honestly + the strategic frame)
# Welsh: Don't hide the price. Frame it.
# Original deck ended with "Either solution will meet BNY's needs" — a thud.
# We turn it into a strategic choice with numbers.
# ----------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, WHITE)
add_accent_bar_top(slide)

add_title_block(slide, "The cost of being right vs. being cheap", eyebrow_text="Why  ·  What it means for BNY")

# Two side-by-side blocks
left_x = Inches(0.8)
right_x = Inches(7.0)
block_w = Inches(5.8)
block_y = Inches(2.0)
block_h = Inches(4.2)

# LEFT: Meraki
left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, block_y, block_w, block_h)
left.fill.solid()
left.fill.fore_color.rgb = LIGHT_GRAY_BG
left.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

add_textbox(slide, left_x + Inches(0.4), block_y + Inches(0.3), block_w - Inches(0.8), Inches(0.4),
            "MERAKI", font_size=Pt(12), color=MID_GRAY, bold=True)
add_textbox(slide, left_x + Inches(0.4), block_y + Inches(0.65), block_w - Inches(0.8), Inches(0.8),
            "$1.55M", font_size=Pt(44), color=WWT_DARK_BLUE, bold=True)
add_textbox(slide, left_x + Inches(0.4), block_y + Inches(1.55), block_w - Inches(0.8), Inches(0.5),
            "Saves $1.4M. Ships in 66 days.",
            font_size=Pt(15), color=BODY_GRAY, bold=True)

# What you give up
add_textbox(slide, left_x + Inches(0.4), block_y + Inches(2.2), block_w - Inches(0.8), Inches(0.35),
            "WHAT YOU GIVE UP", font_size=Pt(11), color=ORANGE, bold=True)

trade_offs = [
    "Two networks to operate, not one",
    "Re-platform when the broader Refresh lands",
    "End-to-end segmentation never reaches the branch"
]
for i, t in enumerate(trade_offs):
    add_textbox(slide, left_x + Inches(0.4), block_y + Inches(2.55) + Inches(0.4)*i,
                block_w - Inches(0.8), Inches(0.4),
                f"·  {t}", font_size=Pt(13), color=BODY_GRAY)

# RIGHT: Cisco SD-WAN
right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, block_y, block_w, block_h)
right.fill.solid()
right.fill.fore_color.rgb = NAVY
right.line.fill.background()

add_textbox(slide, right_x + Inches(0.4), block_y + Inches(0.3), block_w - Inches(0.8), Inches(0.4),
            "CISCO SD-WAN  ·  RECOMMENDED", font_size=Pt(12), color=LIGHT_BLUE_75, bold=True)
add_textbox(slide, right_x + Inches(0.4), block_y + Inches(0.65), block_w - Inches(0.8), Inches(0.8),
            "$2.94M", font_size=Pt(44), color=WHITE, bold=True)
add_textbox(slide, right_x + Inches(0.4), block_y + Inches(1.55), block_w - Inches(0.8), Inches(0.5),
            "$1.4M more today. Zero rework tomorrow.",
            font_size=Pt(15), color=WHITE, bold=True)

# What you get
add_textbox(slide, right_x + Inches(0.4), block_y + Inches(2.2), block_w - Inches(0.8), Inches(0.35),
            "WHAT YOU GET", font_size=Pt(11), color=ORANGE, bold=True)

gets = [
    "One architecture, branch to data center",
    "The pilot pattern for the 2027 rollout volume",
    "Identity-based segmentation, end to end"
]
for i, g in enumerate(gets):
    add_textbox(slide, right_x + Inches(0.4), block_y + Inches(2.55) + Inches(0.4)*i,
                block_w - Inches(0.8), Inches(0.4),
                f"·  {g}", font_size=Pt(13), color=WHITE)

# Bottom takeaway
add_textbox(slide, CONTENT_LEFT, Inches(6.5), CONTENT_WIDTH, Inches(0.5),
            "Meraki saves money on these eleven sites. Cisco SD-WAN saves money on the next two hundred.",
            font_size=Pt(15), color=WWT_DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER, italic=True)

add_monogram(slide)
add_gradient_line(slide)

add_speaker_notes(slide, """WHY — What it means for BNY. Where the deck either wins or loses.

BME: Begin = honest about the cost gap (don't hide it). Middle = the trade is architecture, not features. End = the bottom line — "Meraki saves money on these eleven sites. Cisco saves money on the next two hundred."

What to say (90 seconds):
"Cisco SD-WAN costs $1.4 million more. I want to be honest about that — it's a real number. Here's what that money buys."

Walk through the right column. Land hard on the bottom line at the end.

CRITICAL: Do NOT apologize for the price. Welsh's "Power of I Don't Know" applies in reverse here — the power of "I'm confident" matters. You believe this is the right call. Show it.

REPETITION WITH VARIATION — three ways to land the bottom line:
1. "Pay now or pay twice." (financial)
2. "We're not buying eleven branches. We're buying the pattern for two hundred." (strategic)
3. "The cheaper option is more expensive in eighteen months." (callback to slide 2!)

CURVEBALL PREP: If the COO says "but $1.4M is real money" — agree, then use option 2. If they say "can we do Meraki for some sites and Cisco for others" — that's the trap. Two networks is exactly what slide 2 said broke last time.""")

# ----------------------------------------------------------------------------
# SLIDE 6 — MAGIC WAND (End with wanting — the ask)
# Welsh: Every deck must end with wanting. Original ended with "Either solution
# will meet BNY's needs" — the opposite. We end with what success looks like
# and a clear ask.
# ----------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
set_solid_bg(slide, WHITE)
add_accent_bar_top(slide)

add_title_block(slide, "What we need from you today", eyebrow_text="Magic wand  ·  The ask")

# Picture of success
add_textbox(slide, CONTENT_LEFT, Inches(2.0), CONTENT_WIDTH, Inches(0.45),
            "BY END OF Q3 2026 — IF WE DECIDE THIS WEEK",
            font_size=Pt(12), color=ORANGE, bold=True)

add_textbox(slide, CONTENT_LEFT, Inches(2.5), CONTENT_WIDTH, Inches(0.7),
            "11 microbranches live. 4 head-ends operational. One architecture.",
            font_size=Pt(28), color=WWT_DARK_BLUE, bold=True)

# Three commitments (Rule of 3)
commit_y = Inches(3.7)
commit_w = Inches(3.85)
commits = [
    ("This week", "COO sign-off on Cisco SD-WAN as the RAP architecture.",
     "Unblocks the PO and the 182-day SD-WAN lead time."),
    ("By May 1", "Executed services contract with WWT.",
     "Triggers certification testing and first-article build."),
    ("By Q3 close", "First 11 sites turnkey. WWT BiB pilot complete.",
     "Pattern is proven. 2027 rollout volume can scale.")
]

for i, (when, what, why) in enumerate(commits):
    x = CONTENT_LEFT + (commit_w + Inches(0.18)) * i

    # Top blue bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, commit_y, commit_w, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WWT_LIGHT_BLUE
    bar.line.fill.background()

    add_textbox(slide, x + Inches(0.2), commit_y, commit_w - Inches(0.4), Inches(0.5),
                when.upper(), font_size=Pt(14), color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    # Card body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, commit_y + Inches(0.5), commit_w, Inches(2.2))
    body.fill.solid()
    body.fill.fore_color.rgb = LIGHT_GRAY_BG
    body.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    body.line.width = Pt(0.5)

    add_textbox(slide, x + Inches(0.2), commit_y + Inches(0.7), commit_w - Inches(0.4), Inches(1.0),
                what, font_size=Pt(14), color=WWT_DARK_BLUE, bold=True)

    add_textbox(slide, x + Inches(0.2), commit_y + Inches(1.6), commit_w - Inches(0.4), Inches(1.0),
                why, font_size=Pt(12), color=BODY_GRAY, italic=True)

# Closing line — the wanting
add_textbox(slide, CONTENT_LEFT, Inches(6.6), CONTENT_WIDTH, Inches(0.5),
            "One decision today. Eleven sites by September. The pattern for everything that follows.",
            font_size=Pt(16), color=WWT_DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_monogram(slide)
add_gradient_line(slide)

add_speaker_notes(slide, """MAGIC WAND — End with wanting.

BME: Begin = paint the picture of success. Middle = three concrete commitments with dates. End = the closing line is the message they carry out of the room.

What to say (60 seconds):
"Here's what I need from you today, and here's what we deliver if we get it. Three asks. First one is yours — sign-off this week so we can cut the PO. The other two are ours."

Walk the three cards. Then PAUSE.

Final words (memorize these — say nothing else):
"One decision today. Eleven sites by September. The pattern for everything that follows. — What questions do you have?"

Then SHUT UP. Three-second rule. Make them speak first.

POWER OF "I DON'T KNOW" — REMINDER:
When questions come, you have permission to say "I don't have that detail at my fingertips — let me come back to you with the right answer by end of day." That answer is more powerful than a guess.

REPETITION WITH VARIATION — the closing line in 3 forms (use the one that fits the room):
1. "One decision, eleven sites, one pattern." (rhythmic)
2. "Sign this week, ship by September." (action-forward)
3. "Today's decision is what gets reused two hundred times." (strategic)""")

# ============================================================================
# Save
# ============================================================================
output_path = "/home/claude/output/BNY_RAP_reworked_demo.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
